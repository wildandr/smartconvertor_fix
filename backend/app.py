from flask_cors import CORS
from flask import Flask, request, jsonify
import os
from datetime import datetime
from PyPDF2 import PdfReader
from pptx import Presentation
import nltk
from nltk.tokenize import sent_tokenize, word_tokenize
from nltk import pos_tag
from random import choice
import matplotlib.pyplot as plt
from io import BytesIO
import base64
from PIL import Image
from pptx import Presentation
from io import BytesIO
import base64
from PIL import Image
from pptx import Presentation
from io import BytesIO
import base64
from PIL import Image, ImageDraw
from pptx import Presentation
from pdf2image import convert_from_path
import tempfile
import os
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from subprocess import run
import shutil
import glob
from subprocess import run
import glob
import logging
import json
import sqlite3

nltk.download('punkt')
nltk.download('averaged_perceptron_tagger')

app = Flask(__name__)
CORS(app)

@app.route('/upload', methods=['POST'])
def upload_file():
    if 'file' not in request.files:
        return jsonify({"error": "No file uploaded"}), 400

    file = request.files['file']
    original_filename = file.filename

    timestamp = datetime.now().strftime('%Y-%m-%d_%H-%M-%S')

    preview_image = None

    # Tentukan direktori berdasarkan jenis file
    if original_filename.endswith('.pdf'):
        new_filename = f"{timestamp}_{original_filename}"
        file_path = os.path.join('data_pdf', new_filename)
        file.save(file_path)
        
        # Convert PDF to text and get a preview image
        txt_file_path = convert_pdf_to_txt(file_path)
        preview_image = generate_pdf_preview(file_path)

    elif original_filename.endswith('.ppt') or original_filename.endswith('.pptx'):
        new_filename = f"{timestamp}_{original_filename}"
        file_path = os.path.join('data_ppt', new_filename)
        file.save(file_path)
        
        # Convert PPT to text and get a preview image
        txt_file_path = convert_ppt_to_txt(file_path)
        preview_image = generate_ppt_preview(file_path)

    else:
        return jsonify({"error": "Unsupported file format"}), 400

    # Membuat pertanyaan dan jawaban
    questions = generate_questions(txt_file_path)

    # Menyimpan pertanyaan dan jawaban dalam format JSON
    qa_file = f'json_data/{timestamp}_qa.json'
    with open(qa_file, 'w') as f:
        json.dump(questions, f)

    return jsonify({
        "message": "File processed successfully",
        "preview_image": preview_image,
        "questions": questions
    }), 200

def get_latest_txt_filename():
    list_of_txt_files = glob.glob('data_txt/*.txt')
    latest_file = max(list_of_txt_files, key=os.path.getctime, default=None)

    if latest_file:
        filename = os.path.basename(latest_file)
        filename = os.path.splitext(filename)[0]  # Hilangkan ekstensi file
        filename = filename.split(' ', 1)[-1]  # Hilangkan tanggal dari nama file
        filename = filename.replace('_', ' ')  # Ganti tanda _ dengan spasi
        return filename
    else:
        return None


@app.route('/get_questions', methods=['GET'])
def get_questions():
    # Ambil nama file terbaru dari direktori data_txt
    source_file = get_latest_txt_filename()

    if source_file is None:
        return jsonify({"error": "No text file available"}), 404

    # Temukan file JSON terbaru
    list_of_files = glob.glob('json_data/*.json')
    latest_file = max(list_of_files, key=os.path.getctime, default=None)
    
    if latest_file is None:
        return jsonify({"error": "No JSON data available"}), 404

    with open(latest_file, 'r') as f:
        questions = json.load(f)

    return jsonify({
        "questions": questions,
        "source_file": source_file
    }), 200
    
@app.route('/submit_answer', methods=['POST'])
def submit_answer():
    data = request.get_json()
    answers = data.get('answers')  # Daftar jawaban pengguna

    if not answers:
        return jsonify({"error": "Answers are required"}), 400

    # Membaca jawaban yang benar dari file JSON terakhir
    list_of_files = glob.glob('json_data/*_qa.json')
    latest_file = max(list_of_files, key=os.path.getctime, default=None)

    if latest_file is None:
        return jsonify({"error": "No questions and answers data available"}), 404

    with open(latest_file, 'r') as f:
        questions = json.load(f)

    # Mengumpulkan jawaban yang benar
    correct_answers = {q['question']: q['answer'] for q in questions}

    # Menyiapkan feedback jawaban
    feedback = []
    for answer in answers:
        question_text = answer['question']
        user_answer = answer['answer']
        correct_answer = correct_answers.get(question_text, "No correct answer found")
        
        feedback.append({
            "question": question_text,
            "user_answer": user_answer,
            "correct_answer": correct_answer,
            "is_correct": user_answer == correct_answer
        })

    # Menyimpan feedback ke JSON untuk audit atau review lebih lanjut
    timestamp = datetime.now().strftime('%Y-%m-%d_%H-%M-%S')
    feedback_file = f'json_data/{timestamp}_feedback.json'
    with open(feedback_file, 'w') as f:
        json.dump(feedback, f)

    return jsonify(feedback), 200

@app.route('/get_answers', methods=['GET'])
def get_answers():
    # Membaca jawaban pengguna dari file JSON terakhir
    list_of_files = glob.glob('json_data/*_feedback.json')
    latest_file = max(list_of_files, key=os.path.getctime, default=None)

    if latest_file is None:
        return jsonify({"error": "No answers available"}), 404

    with open(latest_file, 'r') as f:
        feedback_data = json.load(f)

    # Membandingkan jawaban pengguna dengan jawaban yang benar
    feedback = []
    for ans in feedback_data:
        question_text = ans['question']
        user_answer = ans['user_answer']
        correct_answer = ans.get('correct_answer', "No correct answer found")

        feedback.append({
            "question": question_text,
            "user_answer": user_answer,
            "correct_answer": correct_answer,
            "is_correct": user_answer == correct_answer
        })

    return jsonify(feedback), 200


def convert_pdf_to_txt(file_path):
    txt_file_path = file_path.replace('data_pdf/', 'data_txt/').replace('.pdf', '.txt')

    lines = []

    reader = PdfReader(file_path)
    for page in reader.pages:
        text = page.extract_text()
        lines.extend(text.split('\n'))

    with open(txt_file_path, 'w') as txt_file:
        current_paragraph = []

        for line in lines:
            line = line.strip()

            if not line:
                continue

            # Cek jika baris berakhir dengan titik
            if line.endswith('.'):
                current_paragraph.append(line)
                txt_file.write(" ".join(current_paragraph) + '\n\n')
                current_paragraph = []
            else:
                current_paragraph.append(line)

    return txt_file_path

def convert_ppt_to_txt(file_path):
    txt_file_path = file_path.replace('data_ppt/', 'data_txt/').replace('.pptx', '.txt').replace('.ppt', '.txt')

    prs = Presentation(file_path)
    with open(txt_file_path, 'w') as txt_file:
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    txt_file.write(shape.text + '\n')

    return txt_file_path

def generate_questions(txt_file_path, num_questions=3):
    # Baca teks dari file TXT
    with open(txt_file_path, 'r') as txt_file:
        text = txt_file.read()

    # Tokenisasi teks menjadi kalimat
    sentences = sent_tokenize(text)
    questions = []

    for _ in range(num_questions):
        # Pilih kalimat secara acak
        sentence = choice(sentences)

        # Tokenisasi kalimat menjadi kata-kata
        words = word_tokenize(sentence)

        # POS tagging untuk menemukan kata yang berarti
        tagged_words = pos_tag(words)

        # Pilih kata yang akan dihapus dari kata benda atau kata sifat
        valid_words = [w for w in tagged_words if w[1].startswith(('NN', 'JJ'))]

        if valid_words:
            word_to_remove = choice(valid_words)[0]

            # Bentuk pertanyaan fill-in-the-blank
            before, after = sentence.split(word_to_remove, 1)
            question = {
                "question": f"{before}______{after}",
                "before": before,
                "after": after,
                "answer": word_to_remove
            }

            questions.append(question)

    return questions

def generate_pdf_preview(file_path):
    from pdf2image import convert_from_path
    # Mengurangi DPI menjadi lebih rendah untuk mengurangi ukuran file
    images = convert_from_path(file_path, first_page=1, last_page=1, dpi=72, fmt='jpeg')
    
    # Simpan gambar ke dalam buffer dengan kualitas yang lebih rendah
    img_buffer = BytesIO()
    images[0].save(img_buffer, format="JPEG", quality=25)  # Mengurangi kualitas lebih jauh untuk kompresi yang lebih besar
    
    # Encoding base64 untuk dikembalikan sebagai string
    encoded_image = base64.b64encode(img_buffer.getvalue()).decode('utf-8')
    
    return encoded_image


def generate_ppt_preview(file_path):
    logging.info("Mengonversi PPT ke PNG...")
    
    ppt_name = os.path.basename(file_path).replace(".pptx", "").replace(".ppt", "")
    output_dir = "smartconvertor/backend/png_ppt"
    
    os.makedirs(output_dir, exist_ok=True)
    
    try:
        run(["soffice", "--headless", "--convert-to", "png", "--outdir", output_dir, file_path], check=True)
        logging.info("Konversi selesai. Mencari file PNG...")
    except Exception as e:
        logging.error(f"Gagal mengkonversi file PPT: {e}")
        return None

    png_files = glob.glob(f"{output_dir}/{ppt_name}*.png")

    if not png_files:
        logging.error("Tidak ada file PNG yang ditemukan.")
        return None

    png_file = png_files[0]
    logging.info(f"File PNG ditemukan: {png_file}")

    try:
        with Image.open(png_file) as img:
            # Convert the image to RGB format (necessary for consistent encoding)
            rgb_img = img.convert('RGB')
            
            # Save the image to a bytes buffer instead of a file
            with open(png_file, 'rb') as image_file:
                encoded_string = base64.b64encode(image_file.read())
                logging.info(f"Gambar telah dikompresi dan di-encode ke base64.")
                return encoded_string.decode('utf-8')
    except Exception as e:
        logging.error(f"Kesalahan saat membaca atau menyimpan file: {str(e)}")
        return None

if __name__ == '__main__':
    app.run(debug=True)